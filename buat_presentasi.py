# -*- coding: utf-8 -*-
"""Materi presentasi analisis konflik agraria Agrinas–KSO — Unit II Harda."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

OUT = r"C:\Users\Patron\Downloads\sawit lagi\Presentasi_Analisis_Konflik_Agraria_Agrinas_KSO_Polda_Riau.pptx"

# Palette — institutional navy (bukan ungu/krem AI-default)
NAVY = RGBColor(0x0F, 0x2A, 0x44)
NAVY2 = RGBColor(0x1A, 0x3A, 0x5C)
STEEL = RGBColor(0x2C, 0x5F, 0x7C)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)  # terracotta-amber institutional, bukan purple
GOLD = RGBColor(0xB8, 0x8A, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF4, 0xF6, 0xF8)
DARK = RGBColor(0x1E, 0x24, 0x2A)
MUTED = RGBColor(0x5A, 0x66, 0x72)
RED = RGBColor(0xB3, 0x2D, 0x2D)
YELLOW = RGBColor(0xC4, 0x8A, 0x14)
GREEN = RGBColor(0x2E, 0x7D, 0x4F)
LIGHT_LINE = RGBColor(0xD0, 0xD7, 0xDE)


def set_run(run, size=14, bold=False, color=DARK, font="Calibri", italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_rect(slide, l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_round(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    # soften corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def textbox(slide, l, t, w, h, text, size=14, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr",
                                         MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font, italic=italic)
    return box


def add_bullets(slide, l, t, w, h, items, size=13, color=DARK, spacing=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        if spacing:
            p.space_after = Pt(6)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=color)
    return box


def footer(slide, page, total, prs):
    # bottom bar
    add_rect(slide, Inches(0), Inches(7.15), prs.slide_width, Inches(0.35), NAVY)
    textbox(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.28),
            "Unit II Harda · Ditreskrimum Polda Riau  |  Internal — Analisis",
            size=9, color=WHITE, font="Calibri")
    textbox(slide, Inches(11.2), Inches(7.18), Inches(1.5), Inches(0.28),
            f"{page}/{total}", size=9, color=WHITE, align=PP_ALIGN.RIGHT)


def header_bar(slide, prs, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)
    add_rect(slide, Inches(0), Inches(0.08), prs.slide_width, Inches(0.95), NAVY)
    textbox(slide, Inches(0.45), Inches(0.22), Inches(12), Inches(0.4),
            title, size=22, bold=True, color=WHITE, font="Calibri")
    if subtitle:
        textbox(slide, Inches(0.45), Inches(0.62), Inches(12), Inches(0.3),
                subtitle, size=11, color=RGBColor(0xB8, 0xC5, 0xD0), font="Calibri")


def kpi_card(slide, l, t, w, h, value, label, accent=ACCENT):
    add_round(slide, l, t, w, h, OFF)
    add_rect(slide, l, t, Inches(0.08), h, accent)
    textbox(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.3), Inches(0.45),
            value, size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    textbox(slide, l + Inches(0.2), t + Inches(0.6), w - Inches(0.3), Inches(0.45),
            label, size=11, color=MUTED, align=PP_ALIGN.LEFT)


def section_card(slide, l, t, w, h, title, body_lines, title_color=NAVY, bar=ACCENT):
    add_round(slide, l, t, w, h, WHITE)
    # border via thin rect overlay look
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_LINE
    shape.line.width = Pt(1)
    try:
        shape.adjustments[0] = 0.06
    except Exception:
        pass
    add_rect(slide, l, t, w, Inches(0.08), bar)
    textbox(slide, l + Inches(0.2), t + Inches(0.2), w - Inches(0.35), Inches(0.35),
            title, size=13, bold=True, color=title_color)
    add_bullets(slide, l + Inches(0.15), t + Inches(0.55), w - Inches(0.3), h - Inches(0.7),
                body_lines, size=11, color=DARK)


def set_cell(cell, text, bold=False, size=10, color=DARK, fill=None, align=PP_ALIGN.CENTER):
    cell.text = ""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    set_run(run, size=size, bold=bold, color=color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    TOTAL = 14

    # ─── 1 COVER ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), prs.slide_height, ACCENT)
    add_rect(s, Inches(0), Inches(6.9), prs.slide_width, Inches(0.6), NAVY2)
    textbox(s, Inches(0.7), Inches(1.1), Inches(11), Inches(0.35),
            "KEPOLISIAN NEGARA REPUBLIK INDONESIA  ·  DAERAH RIAU",
            size=12, color=GOLD, bold=True)
    textbox(s, Inches(0.7), Inches(1.5), Inches(11), Inches(0.3),
            "DIREKTORAT RESERSE KRIMINAL UMUM  —  UNIT II HARDA",
            size=13, color=RGBColor(0xB8, 0xC5, 0xD0))
    textbox(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.9),
            "PAPARAN HASIL ANALISIS",
            size=18, bold=True, color=ACCENT)
    textbox(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.2),
            "Konflik Agraria Agrinas–KSO\ndi 12 Polres Wilayah Hukum Polda Riau",
            size=32, bold=True, color=WHITE)
    textbox(s, Inches(0.7), Inches(4.6), Inches(11), Inches(0.5),
            "Analisis berbasis dokumen Polri / Intelkam / Reskrim  ·  Bukan verifikasi lapangan",
            size=13, italic=True, color=RGBColor(0x9A, 0xA8, 0xB5))
    textbox(s, Inches(0.7), Inches(5.5), Inches(8), Inches(0.7),
            "Periode sumber: 2024 – Juli 2026\nTanggal penyusunan: 3 Agustus 2026",
            size=12, color=WHITE)
    textbox(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.35),
            "RAHASIA INTERNAL — UNTUK KEPENTINGAN ANALISIS UNIT II HARDA",
            size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ─── 2 AGENDA ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "Alur Paparan", "Enam bagian utama")
    agenda = [
        ("01", "Ringkasan Eksekutif", "Temuan inti & skala ekosistem 130 perusahaan"),
        ("02", "Profil Tingkat Polda", "Klaster Intelkam, bentrok, ranking prioritas"),
        ("03", "Peta Panas Koridor", "Utara · Selatan · Tengah · Nihil Agrinas"),
        ("04", "Titik Merah & Hotspot", "Rohul · Rohil · Bengkalis + satker kritis"),
        ("05", "Sintesis Analitik", "Rantai risiko, tipologi, diskoneksi data"),
        ("06", "Rekomendasi", "Agenda riset R1–R8 & lima prioritas segera"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        col = i % 3
        row = i // 3
        l = Inches(0.5 + col * 4.2)
        t = Inches(1.4 + row * 2.5)
        add_round(s, l, t, Inches(3.9), Inches(2.2), OFF)
        textbox(s, l + Inches(0.25), t + Inches(0.3), Inches(3.4), Inches(0.45),
                num, size=28, bold=True, color=ACCENT)
        textbox(s, l + Inches(0.25), t + Inches(0.9), Inches(3.4), Inches(0.4),
                title, size=16, bold=True, color=NAVY)
        textbox(s, l + Inches(0.25), t + Inches(1.35), Inches(3.4), Inches(0.6),
                desc, size=12, color=MUTED)
    footer(s, 2, TOTAL, prs)

    # ─── 3 RINGKASAN EKSEKUTIF ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "1. Ringkasan Eksekutif", "Temuan inti untuk Unit II Harda")
    kpi_card(s, Inches(0.4), Inches(1.3), Inches(2.4), Inches(1.15), "130", "Perusahaan Agrinas–KSO", STEEL)
    kpi_card(s, Inches(3.0), Inches(1.3), Inches(2.4), Inches(1.15), "3 · 24 · 103", "Merah · Kuning · Hijau", RED)
    kpi_card(s, Inches(5.6), Inches(1.3), Inches(2.4), Inches(1.15), "1 MD + 23", "Korban jiwa & luka", ACCENT)
    kpi_card(s, Inches(8.2), Inches(1.3), Inches(2.4), Inches(1.15), "6", "Bentrok (utara)", YELLOW)
    kpi_card(s, Inches(10.8), Inches(1.3), Inches(2.1), Inches(1.15), "Inhu", "Volume kebun tertinggi", GREEN)

    findings = [
        "Konflik struktural tata kelola (penunjukan KSO, PAM non-BUJP, batas pasca-PKH, plasma/MHA) lebih determinan daripada pencurian TBS oportunistik.",
        "Tiga titik merah resmi: Gunung Mas Raya–UTS (Rohil), Berkat Satu–Majuma (Rohul, 1 MD), SIS–PAB (Bengkalis).",
        "Pelalawan: KSO seluruhnya hijau, tetapi dimensi TNTN menghasilkan aksi kolektif tinggi — dipisah dari analisis KSO murni.",
        "Pekanbaru & Kep. Meranti hampir NIHIL Agrinas–KSO; jangan digabung dengan pidana pencurian masyarakat.",
        "Kualitas data tidak merata; angka luas Pelalawan / Siak / sebagian Kuansing tidak layak diagregasi tanpa klarifikasi.",
    ]
    add_bullets(s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(4.0), findings, size=14)
    footer(s, 3, TOTAL, prs)

    # ─── 4 SKALA & KLASTER ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "2. Skala Ekosistem & Klaster Intelkam", "Baseline 12 Polres")

    # cluster definition cards
    for i, (lab, desc, col) in enumerate([
        ("MERAH", "Konflik + kerusuhan + korban", RED),
        ("KUNING", "Konflik tanpa kerusuhan", YELLOW),
        ("HIJAU", "Belum terjadi konflik", GREEN),
    ]):
        l = Inches(0.45 + i * 4.25)
        add_round(s, l, Inches(1.25), Inches(4.0), Inches(0.95), OFF)
        add_rect(s, l, Inches(1.25), Inches(0.12), Inches(0.95), col)
        textbox(s, l + Inches(0.3), Inches(1.35), Inches(3.5), Inches(0.35),
                lab, size=16, bold=True, color=col)
        textbox(s, l + Inches(0.3), Inches(1.7), Inches(3.5), Inches(0.35),
                desc, size=12, color=MUTED)

    rows_data = [
        ["Satker", "Kebun", "M", "K", "H", "Konflik"],
        ["Inhu", "28", "0", "2", "26", "3"],
        ["Rohil", "24", "1", "6", "17", "7"],
        ["Rohul", "21", "1", "5", "15", "6"],
        ["Inhil", "15", "0", "4", "11", "4"],
        ["Kuansing", "11", "0", "1", "10", "1"],
        ["Pelalawan", "9", "0", "0", "9", "0"],
        ["Bengkalis", "8", "1", "0", "7", "3"],
        ["Siak", "6", "0", "0", "6", "1"],
        ["Kampar", "4", "0", "3", "1", "3"],
        ["Dumai", "3", "0", "3", "0", "3"],
        ["Meranti", "1", "0", "0", "1", "0"],
        ["TOTAL", "130", "3", "24", "103", "—"],
    ]
    table = s.shapes.add_table(len(rows_data), 6, Inches(0.4), Inches(2.4), Inches(8.2), Inches(4.4)).table
    widths = [1.8, 1.2, 0.9, 0.9, 1.0, 1.4]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            if r == 0:
                set_cell(table.cell(r, c), val, bold=True, size=10, color=WHITE, fill=NAVY)
            elif r == len(rows_data) - 1:
                set_cell(table.cell(r, c), val, bold=True, size=10, color=WHITE, fill=STEEL)
            else:
                fill = OFF if r % 2 else WHITE
                # highlight merah cells
                if c == 2 and val == "1":
                    set_cell(table.cell(r, c), val, bold=True, size=10, color=WHITE, fill=RED)
                else:
                    set_cell(table.cell(r, c), val, size=10, fill=fill)

    # side notes
    add_round(s, Inches(8.9), Inches(2.4), Inches(4.0), Inches(4.4), OFF)
    textbox(s, Inches(9.15), Inches(2.6), Inches(3.5), Inches(0.4),
            "Catatan kunci", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(9.1), Inches(3.15), Inches(3.6), Inches(3.4), [
        "Pekanbaru: NIHIL (di luar 130).",
        "Konflik resume tertinggi: Rohil (7) & Rohul (6).",
        "Volume terbesar: Inhu (28) — bukan intensitas kekerasan tertinggi.",
        "Klaster hijau sering underestimate gesekan lokal.",
    ], size=12)
    footer(s, 4, TOTAL, prs)

    # ─── 5 RANKING PRIORITAS ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "2. Ranking Prioritas Analitik Harda", "Alokasi perhatian Unit II")

    tiers = [
        ("SANGAT TINGGI", "Rohul · Rohil · Bengkalis", "Merah + bentrok + korban jiwa/luka", RED),
        ("TINGGI", "Inhu · Kuansing · Kampar · Dumai · Pelalawan (TNTN)",
         "Volume kebun / LP terstruktur / dimensi kawasan", ACCENT),
        ("SEDANG", "Inhil · Siak", "Konflik ada; data masih parsial", YELLOW),
        ("RENDAH / GAP-FILL", "Kep. Meranti · Pekanbaru", "NIHIL Agrinas–KSO; monitoring ringan", GREEN),
    ]
    for i, (tier, satker, alasan, col) in enumerate(tiers):
        t = Inches(1.3 + i * 1.35)
        add_round(s, Inches(0.5), t, Inches(12.3), Inches(1.2), OFF)
        add_rect(s, Inches(0.5), t, Inches(0.15), Inches(1.2), col)
        textbox(s, Inches(0.9), t + Inches(0.2), Inches(3.2), Inches(0.35),
                tier, size=14, bold=True, color=col)
        textbox(s, Inches(4.2), t + Inches(0.2), Inches(8.2), Inches(0.35),
                satker, size=15, bold=True, color=NAVY)
        textbox(s, Inches(4.2), t + Inches(0.6), Inches(8.2), Inches(0.4),
                alasan, size=13, color=MUTED)
    footer(s, 5, TOTAL, prs)

    # ─── 6 PETA PANAS KORIDOR ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "3. Peta Panas Komparatif 12 Polres", "Empat koridor analitik")

    corridors = [
        ("UTARA", "Eskalasi fisik", "Rohul · Rohil\nBengkalis · Dumai",
         "Merah/kuning\nPAM · bentrok\nkorban jiwa/luka", RED),
        ("SELATAN", "Volume & tuntutan", "Inhu · Kuansing\nInhil",
         "Banyak kebun\nplasma / MHA\nrelokasi · TBS", STEEL),
        ("TENGAH", "Campuran", "Kampar · Siak\nPelalawan",
         "LP Kampar bersih\nSiak admin/HTI\nPelalawan TNTN", ACCENT),
        ("NIHIL", "Gap-fill", "Pekanbaru\nKep. Meranti",
         "Bukan konflik\ntata kelola\nAgrinas–KSO", GREEN),
    ]
    for i, (name, sub, satker, char, col) in enumerate(corridors):
        l = Inches(0.4 + i * 3.2)
        add_round(s, l, Inches(1.35), Inches(3.0), Inches(5.3), OFF)
        add_rect(s, l, Inches(1.35), Inches(3.0), Inches(0.7), col)
        textbox(s, l + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.5),
                name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, l + Inches(0.2), Inches(2.25), Inches(2.6), Inches(0.35),
                sub, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        textbox(s, l + Inches(0.2), Inches(2.8), Inches(2.6), Inches(1.3),
                satker, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        textbox(s, l + Inches(0.2), Inches(4.3), Inches(2.6), Inches(1.8),
                char, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 6, TOTAL, prs)

    # ─── 7 TITIK MERAH ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "4. Tiga Titik Merah Resmi", "Koridor utara — dampak tertinggi")

    merah = [
        ("ROHUL", "Berkat Satu – Majuma",
         ["KSO Nusantara Sawit Majuma (Bonai)", "1 meninggal dunia + luka",
          "LP/B/07/II/2026 — P21", "Bentrok PAM swakarsa vs PAM"], RED),
        ("ROHIL", "Gunung Mas Raya – UTS",
         ["KSO PT Ujung Tanjung Sejahtera", "PAM Flores non-BUJP",
          "Bentrok 20 Okt 2025 — 7 luka", "Penolakan serah/plang KSO"], ACCENT),
        ("BENGKALIS", "SIS – PAB (Mandau)",
         ["Bentrok berulang Des 2025–Mei 2026", "Multi-LP; kerusakan kendaraan",
          "Panen liar ~150 Ha; bakar pos", "Klaster resume vs kartu tidak selaras"], YELLOW),
    ]
    for i, (wil, title, points, col) in enumerate(merah):
        l = Inches(0.4 + i * 4.25)
        add_round(s, l, Inches(1.3), Inches(4.05), Inches(5.4), OFF)
        add_rect(s, l, Inches(1.3), Inches(4.05), Inches(0.9), col)
        textbox(s, l + Inches(0.2), Inches(1.4), Inches(3.6), Inches(0.3),
                wil, size=12, bold=True, color=WHITE)
        textbox(s, l + Inches(0.2), Inches(1.7), Inches(3.6), Inches(0.35),
                title, size=14, bold=True, color=WHITE)
        add_bullets(s, l + Inches(0.2), Inches(2.5), Inches(3.6), Inches(3.8), points, size=13)
    footer(s, 7, TOTAL, prs)

    # ─── 8 HOTSPOT LAIN ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "4. Hotspot & Satker Kritis Lain", "Di luar tiga titik merah resmi")

    cards = [
        ("INHU — Volume tertinggi", [
            "28 kebun; 0M / 2K / 26H (underestimate)",
            "Plasma, MHA Talang Mamak, panen paksa",
            "EX Palm 560 Ha (Apr 2026); Bernas Mulya Mandiri",
        ]),
        ("KAMPAR — LP terstruktur", [
            "Intelkam 4 kebun vs lokal 21 estate",
            "Agrinas 29.174 Ha; LP Jun 2026 Pasir Sialang",
            "Pisahkan LP Agrinas vs agraria umum",
        ]),
        ("KUANSING — TBS berulang", [
            "~35–37 LP relevan Agrinas-estate",
            "Cerenti, WJT, Pesikaian (relokasi TNTN)",
            "Pisahkan pidana oportunistik vs struktural",
        ]),
        ("PELALAWAN — TNTN", [
            "9 KSO hijau; konflik KSO resume 0",
            "TNTN ~81.793 Ha; 7 desa; 5.733 KK",
            "Portofolio risiko terpisah dari KSO",
        ]),
        ("DUMAI — Semua kuning", [
            "DMMP–Riden Jaya (Feb 2026)",
            "Perebutan pekerjaan SRPO & Pelintung",
            "Risiko double-count dengan Bengkalis",
        ]),
        ("MUTIARA NAGA — Watchlist", [
            "Bengkalis: mediasi gagal 28 Jul 2026",
            "Bentrok 29 Jul — kandidat naik klaster",
            "Validasi klaster Intelkam mendesak",
        ]),
    ]
    for i, (title, pts) in enumerate(cards):
        col = i % 3
        row = i // 3
        l = Inches(0.4 + col * 4.25)
        t = Inches(1.25 + row * 2.85)
        add_round(s, l, t, Inches(4.05), Inches(2.65), OFF)
        add_rect(s, l, t, Inches(4.05), Inches(0.08), ACCENT if i < 5 else RED)
        textbox(s, l + Inches(0.2), t + Inches(0.25), Inches(3.6), Inches(0.4),
                title, size=13, bold=True, color=NAVY)
        add_bullets(s, l + Inches(0.15), t + Inches(0.75), Inches(3.7), Inches(1.7), pts, size=12)
    footer(s, 8, TOTAL, prs)

    # ─── 9 RANTAI RISIKO ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "5. Rantai Tata Kelola Berisiko", "Pola generik eskalasi")

    steps = [
        ("1", "Sita\nSatgas PKH"),
        ("2", "Titip kelola\nAgrinas"),
        ("3", "Penunjukan\nKSO"),
        ("4", "Non-tempatan\n+ PAM non-BUJP"),
        ("5", "Batas/plasma\nbelum selesai"),
        ("6", "LP / bentrok\n/ MD"),
    ]
    for i, (num, lab) in enumerate(steps):
        l = Inches(0.35 + i * 2.15)
        add_round(s, l, Inches(1.8), Inches(1.9), Inches(2.0), NAVY if i < 5 else RED)
        textbox(s, l + Inches(0.1), Inches(1.95), Inches(1.7), Inches(0.4),
                num, size=20, bold=True, color=ACCENT if i < 5 else WHITE, align=PP_ALIGN.CENTER)
        textbox(s, l + Inches(0.1), Inches(2.45), Inches(1.7), Inches(1.1),
                lab, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 5:
            textbox(s, l + Inches(1.75), Inches(2.4), Inches(0.4), Inches(0.5),
                    "→", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_round(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.3), OFF)
    textbox(s, Inches(0.8), Inches(4.5), Inches(11.8), Inches(0.4),
            "Titik gagal paling berulang", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.8), Inches(5.05), Inches(11.5), Inches(1.4), [
        "Penunjukan KSO tanpa legitimasi sosial lokal.",
        "Pengamanan non-formal (PAM swakarsa non-BUJP) — termasuk isu SARA/premanisme.",
        "Implikasi Harda: prioritas pada rantai penunjukan & pengamanan KSO, bukan semata statistik pencurian TBS.",
    ], size=14)
    footer(s, 9, TOTAL, prs)

    # ─── 10 STRUKTURAL VS OPORTUNISTIK ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "5. Tipologi: Struktural vs Oportunistik", "Pemisahan analitik wajib")

    # left
    add_round(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.4), OFF)
    add_rect(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(0.7), RED)
    textbox(s, Inches(0.6), Inches(1.4), Inches(5.7), Inches(0.5),
            "STRUKTURAL TATA KELOLA", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.0), [
        "Penolakan KSO / plang / serah terima",
        "PAM vs PAM; PAM non-BUJP",
        "Klaim plasma / MHA / desa",
        "TNTN & relokasi kawasan",
        "Contoh: Rohul, Rohil, Bengkalis, Inhu, Pelalawan (TNTN), Kampar (Johan Sentosa)",
        "→ Prioritas pemetaan Harda",
    ], size=14)

    # right
    add_round(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.4), OFF)
    add_rect(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.7), STEEL)
    textbox(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(0.5),
            "PIDANA OPORTUNISTIK", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(7.1), Inches(2.3), Inches(5.5), Inches(4.0), [
        "Pencurian TBS berulang",
        "Tipiring / RJ di areal kebun",
        "Tidak selalu terkait konflik KSO",
        "Contoh: Kuansing, Pekanbaru (non-Agrinas), sebagian Inhu",
        "→ Jangan menggelembungkan “konflik Agrinas” di dashboard",
        "→ Filter tipologi sebelum agregasi",
    ], size=14)
    footer(s, 10, TOTAL, prs)

    # ─── 11 TEMUAN KOMPARATIF + DISKONEKSI ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "5. Temuan Komparatif & Diskoneksi Data", "Implikasi kualitas bukti")

    left_items = [
        "Korban jiwa/luka hampir seluruhnya di koridor utara terkait PAM/KSO.",
        "Jumlah kebun ≠ intensitas kekerasan (Inhu terbanyak; Rohul paling mematikan).",
        "Klaster hijau tidak menjamin aman (WJT, EX Palm, Mutiara Naga).",
        "TNTN Pelalawan = portofolio risiko terpisah.",
        "Double-count geografis harus dibersihkan sebelum skor risiko.",
    ]
    add_round(s, Inches(0.4), Inches(1.25), Inches(6.3), Inches(5.5), OFF)
    textbox(s, Inches(0.65), Inches(1.45), Inches(5.8), Inches(0.4),
            "Temuan komparatif kunci", size=15, bold=True, color=NAVY)
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.4), left_items, size=13)

    layers = [
        ("Intelkam", "Cakupan 130 + klaster", "Underestimate hijau bergesekan"),
        ("Reskrim/LP", "Nomor LP & proses", "Tidak seragam antar satker"),
        ("Harda", "Tipologi yuridis–sosial", "Belum master LP–kebun"),
        ("Disbun/Pemda", "Konteks izin/batas", "Luas sering tidak sinkron"),
    ]
    add_round(s, Inches(6.9), Inches(1.25), Inches(6.0), Inches(5.5), OFF)
    textbox(s, Inches(7.15), Inches(1.45), Inches(5.5), Inches(0.4),
            "Diskoneksi lapisan data", size=15, bold=True, color=NAVY)
    for i, (lap, kuat, lemah) in enumerate(layers):
        t = Inches(2.05 + i * 1.1)
        textbox(s, Inches(7.2), t, Inches(5.4), Inches(0.3),
                lap, size=13, bold=True, color=ACCENT)
        textbox(s, Inches(7.2), t + Inches(0.3), Inches(5.4), Inches(0.55),
                f"+ {kuat}\n− {lemah}", size=11, color=DARK)
    footer(s, 11, TOTAL, prs)

    # ─── 12 REKOMENDASI ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "6. Agenda Analisis Lanjutan (R1–R8)", "Bukan perintah operasi — agenda riset/analitik")

    recs = [
        ("R1", "Matriks risiko satker", "TINGGI", RED),
        ("R2", "Jaringan aktor KSO/PAM", "TINGGI", RED),
        ("R3", "Validasi klaster Intelkam", "TINGGI", RED),
        ("R4", "Timeline eskalasi PKH→KSO→LP", "TINGGI", RED),
        ("R5", "Tipologi hukum Harda", "SEDANG", YELLOW),
        ("R6", "Spasial kecamatan/desa", "SEDANG", YELLOW),
        ("R7", "Gap-fill satker tipis", "SEDANG", YELLOW),
        ("R8", "Dashboard monitoring", "RENDAH", GREEN),
    ]
    for i, (code, title, prio, col) in enumerate(recs):
        col_i = i % 4
        row = i // 4
        l = Inches(0.4 + col_i * 3.2)
        t = Inches(1.35 + row * 2.7)
        add_round(s, l, t, Inches(3.0), Inches(2.4), OFF)
        textbox(s, l + Inches(0.2), t + Inches(0.3), Inches(2.6), Inches(0.4),
                code, size=22, bold=True, color=NAVY)
        textbox(s, l + Inches(0.2), t + Inches(0.9), Inches(2.6), Inches(0.7),
                title, size=14, bold=True, color=DARK)
        # priority pill
        pill = add_round(s, l + Inches(0.2), t + Inches(1.7), Inches(1.6), Inches(0.4), col)
        textbox(s, l + Inches(0.2), t + Inches(1.75), Inches(1.6), Inches(0.35),
                prio, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    footer(s, 12, TOTAL, prs)

    # ─── 13 LIMA PRIORITAS ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, WHITE)
    header_bar(s, prs, "6. Lima Prioritas Segera — Pimpinan Unit II", "Urutan kerja analitik yang disarankan")

    prioritas = [
        ("1", "R1 — Skor risiko satker", "Ranking bulanan berbasis kebun, klaster, LP, korban, PAM non-BUJP"),
        ("2", "R3 — Validasi klaster", "Utamanya Bengkalis & kebun hijau bergesekan (Mutiara Naga, WJT, EX Palm, Johan Sentosa)"),
        ("3", "R2 — Jaringan KSO/PAM", "Seed: Bernas Mulya, Runggu, PAB, Majuma, Digjaya/UTS, Riden Jaya"),
        ("4", "R4 — Timeline eskalasi", "Pilot: Majuma (Rohul), PAB–SIS (Bengkalis), UTS (Rohil), DMMP (Dumai)"),
        ("5", "R7 — Gap-fill data", "Rohil LP + Pelalawan daftar KSO + Meranti template satker"),
    ]
    for i, (num, title, desc) in enumerate(prioritas):
        t = Inches(1.25 + i * 1.1)
        add_round(s, Inches(0.5), t, Inches(12.3), Inches(1.0), OFF)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), t + Inches(0.2), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if i < 3 else STEEL
        circle.line.fill.background()
        textbox(s, Inches(0.7), t + Inches(0.3), Inches(0.6), Inches(0.45),
                num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        textbox(s, Inches(1.6), t + Inches(0.15), Inches(10.8), Inches(0.35),
                title, size=16, bold=True, color=NAVY)
        textbox(s, Inches(1.6), t + Inches(0.52), Inches(10.8), Inches(0.35),
                desc, size=12, color=MUTED)
    footer(s, 13, TOTAL, prs)

    # ─── 14 PENUTUP ───
    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), prs.slide_height, ACCENT)
    textbox(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
            "PENUTUP", size=16, bold=True, color=ACCENT)
    textbox(s, Inches(0.8), Inches(2.1), Inches(11.5), Inches(1.5),
            "Konflik Agrinas–KSO di Riau adalah persoalan\ntata kelola dan pengamanan — bukan semata pencurian TBS.",
            size=24, bold=True, color=WHITE)
    add_bullets(s, Inches(0.8), Inches(4.0), Inches(11), Inches(1.8), [
        "Fokus analitik: rantai penunjukan KSO + PAM non-BUJP + validasi klaster.",
        "Produk pendukung: matriks_agrinas_kso_12_polres.xlsx & laporan Word lengkap.",
        "Sifat: deskripsi–sintesis dokumen; bukan verifikasi lapangan / perintah operasi.",
    ], size=14, color=RGBColor(0xC8, 0xD2, 0xDA))
    textbox(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
            "Unit II Harda  ·  Ditreskrimum Polda Riau  ·  Agustus 2026",
            size=13, color=GOLD)

    prs.save(OUT)
    print(f"OK: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
